"""
PPTX Scrubber — anonymize PowerPoint presentations.

Walks all slides, shapes, text frames, and tables.
Preserves formatting (font, size, colour, bold/italic) while
replacing only the text content of each run.

Requires: python-pptx
"""

import os
from typing import Dict, List, Tuple

from .preset import Preset
from .layer1 import analyze_layer1_text
from .findings import Finding



def scrub_pptx(
    input_path: str,
    output_path: str,
    preset: Preset,
    language: str,
    file_id: str,
) -> Tuple[List[Finding], Dict[str, int]]:
    """
    Anonymize a PowerPoint file.

    Args:
        input_path:   Path to the source .pptx file.
        output_path:  Destination path for the redacted .pptx.
        preset:       Anonymization configuration.
        language:     BCP-47 language code, e.g. "en", "de".
        file_id:      Identifier used in Finding records.

    Returns:
        (findings, summary) — same contract as docx_scrubber.scrub_docx.
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ImportError(
            "python-pptx is required for PPTX support. "
            "Install it with: pip install python-pptx"
        ) from exc

    # Select analysis function based on preset layer
    if preset.layer == 2:
        from .layer2_candle import analyze_layer2_text as _analyze
    elif preset.layer == 3:
        from .layer3_presidio import analyze_layer3_text as _analyze
    else:
        _analyze = analyze_layer1_text

    prs = Presentation(input_path)
    all_findings: List[Finding] = []
    summary: Dict[str, int] = {}
    filename = os.path.basename(input_path)

    def process(text: str, location: str) -> str:
        """Redact text and collect findings."""
        redacted, findings, local = _analyze(text, preset, language)
        for f in findings:
            f.file_id = file_id
            f.original_filename = filename
            f.page_or_location = location
        all_findings.extend(findings)
        for k, v in local.items():
            summary[k] = summary.get(k, 0) + v
        return redacted

    def scrub_text_frame(tf, location: str) -> None:
        """Process all runs in a text frame in place."""
        for para in tf.paragraphs:
            for run in para.runs:
                if run.text and run.text.strip():
                    run.text = process(run.text, location)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_label = f"slide {slide_idx}"

        for shape_idx, shape in enumerate(slide.shapes, start=1):
            shape_label = f"{slide_label}, shape {shape_idx}"

            # Text frames (titles, text boxes, placeholders)
            if shape.has_text_frame:
                scrub_text_frame(shape.text_frame, shape_label)

            # Tables
            if shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows, start=1):
                    for col_idx, cell in enumerate(row.cells, start=1):
                        cell_label = f"{slide_label}, table row {row_idx} col {col_idx}"
                        scrub_text_frame(cell.text_frame, cell_label)

    prs.save(output_path)
    return all_findings, summary
